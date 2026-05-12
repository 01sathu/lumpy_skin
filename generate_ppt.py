from pptx import Presentation
from pptx.util import Inches, Pt
import os

def create_presentation():
    prs = Presentation()

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "AI-Powered Bovine Disease Diagnostic Center"
    subtitle.text = "Automated Detection of Lumpy Skin Disease (LSD) & Foot and Mouth Disease (FMD)\nPresenter: Lead Researcher"

    # Slide 2: The Problem
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Why This Matters?"
    tf = body_shape.text_frame
    tf.text = "Cattle diseases like Lumpy Skin Disease (LSD) and Foot and Mouth Disease (FMD) spread rapidly."
    p = tf.add_paragraph()
    p.text = "Late diagnosis leads to severe economic losses for farmers and the dairy industry."
    p = tf.add_paragraph()
    p.text = "Access to immediate veterinary expertise is often limited in rural areas."

    # Slide 3: Our Solution
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Our Solution: An AI-Powered Diagnostic Assistant"
    tf = body_shape.text_frame
    tf.text = "Instant Analysis: A web-based tool that diagnoses diseases from a simple photo."
    p = tf.add_paragraph()
    p.text = "Highly Accurate: Powered by state-of-the-art Deep Learning (AI)."
    p = tf.add_paragraph()
    p.text = "Accessible Anywhere: Can be used on any smartphone or computer with internet access."

    # Slide 4: Data
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Data Collection & Preparation"
    tf = body_shape.text_frame
    tf.text = "Gathering Data: Collected over 1,000+ real-world images of LSD, FMD, and healthy cattle."
    p = tf.add_paragraph()
    p.text = "Ensuring Quality: Filtered for medical accuracy and diverse lighting conditions."
    p = tf.add_paragraph()
    p.text = "Training vs. Testing: Split data (80% for training the AI, 20% for testing its accuracy)."

    # Slide 5: The Technology
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "The Technology: YOLOv8 AI Model"
    tf = body_shape.text_frame
    tf.text = "What is it? YOLO (You Only Look Once) is an advanced vision AI."
    p = tf.add_paragraph()
    p.text = "Why YOLOv8? It is extremely fast and highly accurate for image classification."
    p = tf.add_paragraph()
    p.text = "The Process: The AI learned to recognize the specific skin nodules of LSD and the lesions of FMD."

    # Slide 6: Results
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "The Results: High Accuracy & Reliability"
    tf = body_shape.text_frame
    tf.text = "Target Accuracy: > 95%"
    p = tf.add_paragraph()
    p.text = "Achieved Accuracy: ~ 95% on our strict research-grade dataset."
    p = tf.add_paragraph()
    p.text = "Minimal Errors: The model rarely confuses a diseased cow with a healthy one."

    # Slide 7: Web Application
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "The Web Application"
    tf = body_shape.text_frame
    tf.text = "Built using Streamlit for a clean, user-friendly interface."
    p = tf.add_paragraph()
    p.text = "Features:"
    p.level = 0
    p1 = tf.add_paragraph()
    p1.text = "Dedicated portals for LSD and FMD."
    p1.level = 1
    p2 = tf.add_paragraph()
    p2.text = "Instant image upload and prediction."
    p2.level = 1
    p3 = tf.add_paragraph()
    p3.text = "Easy-to-read confidence scores."
    p3.level = 1

    # Slide 8: Future Roadmap
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "What's Next? (Future Roadmap)"
    tf = body_shape.text_frame
    tf.text = "Mobile App: Creating a dedicated offline mobile app for remote villages."
    p = tf.add_paragraph()
    p.text = "More Diseases: Expanding the AI to detect other issues like Mastitis."
    p = tf.add_paragraph()
    p.text = "Precision Targeting: Using object detection to highlight the exact location of lesions on the screen."

    # Slide 9: Conclusion
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Conclusion"
    tf = body_shape.text_frame
    tf.text = "Early detection saves livestock and livelihoods."
    p = tf.add_paragraph()
    p.text = "AI is a powerful tool to support, not replace, veterinary professionals."
    p = tf.add_paragraph()
    p.text = "Thank You! Questions?"

    ppt_path = "Bovine_Disease_Diagnostic_Center_Presentation.pptx"
    prs.save(ppt_path)
    print(f"Presentation saved to {ppt_path}")

if __name__ == '__main__':
    create_presentation()
